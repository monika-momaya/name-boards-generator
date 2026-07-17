"""
Name Board Generator — core engine.

Generates A4-landscape PPTX where each slide contains a fold-over tent card:
  - Top half: NAME + TITLE/COMPANY, rotated 180°
  - Bottom half: NAME + TITLE/COMPANY, upright

Fonts:
  - Name:           AlternateGothic2 BT, bold  (ALL CAPS, 90pt)
  - Title/Company:  AlternateGothic2 BT (Title Case, allowed to wrap to 2 lines)

Layout rule: if Title and Company each fit on their own line (within the
max line count budget), they are stacked tightly (own lines). If either
would overflow its line budget, Title and Company are merged into a single
comma-separated line instead.
"""

from __future__ import annotations

import copy
import io
import uuid
import zipfile
from dataclasses import dataclass
from typing import Optional

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import ImageFont
import os

# ---------------------------------------------------------------------------
# Constants
# ---------------------------------------------------------------------------

# A4 landscape
SLIDE_W_IN = 11.69
SLIDE_H_IN = 8.27

# These are the actual intended fonts. PowerPoint will use them correctly
# on any machine that has them installed; otherwise it substitutes a
# default font. Upload a matching .ttf/.otf in the app sidebar to also get
# accurate auto-shrink/wrap measurements during generation.
FONT_NAME_BOLD = "AlternateGothic2 BT"
FONT_NAME_MEDIUM = "AlternateGothic2 BT"

# Fallback fonts used ONLY for measuring text width if the real TTF files
# are not available to PIL/Pillow in this environment. PowerPoint itself
# will always be told to use FONT_NAME_BOLD / FONT_NAME_MEDIUM regardless,
# so visual output in real PowerPoint is correct once fonts are installed
# there. This fallback only affects the auto-shrink *measurement* step.
FALLBACK_FONT_REGULAR = "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf"

NAME_COLOR = RGBColor(0x00, 0x00, 0x00)
TITLE_COLOR = RGBColor(0x00, 0x00, 0x00)

# Box geometry
# Textbox width is fixed at 29 cm; margin is derived to centre it on the slide.
# SLIDE_W_IN = 11.69 in = 29.693 cm  →  each side margin = (29.693 - 29) / 2 cm
TEXTBOX_W_IN = 29 / 2.54          # 29 cm → inches
MARGIN_X     = (SLIDE_W_IN - TEXTBOX_W_IN) / 2   # ≈ 0.136 in per side

# Fixed font sizes. MAX == MIN disables auto-shrink; the PPTX is fully
# editable so any overflow can be resized by hand in PowerPoint.
NAME_MAX_PT  = 95
NAME_MIN_PT  = 95
TITLE_MAX_PT = 55
TITLE_MIN_PT = 55

# Vertical gap between Name block and Title block (loose)
NAME_TITLE_GAP_IN = 0.18
# Vertical gap between Title line and Company line (tight)
TITLE_COMPANY_GAP_IN = 0.02

HALF_H_IN = SLIDE_H_IN / 2

# Registered custom font paths (filled in by register_fonts())
# If a licensed font file has been dropped locally into fonts/, pick it up
# automatically for measurement. Nothing ships here by default (see README).
_DEFAULT_FONT_PATH = os.path.join(os.path.dirname(__file__), "fonts", "ALTGOT2N.TTF")
_CUSTOM_FONT_PATHS: dict[str, str] = {}
if os.path.isfile(_DEFAULT_FONT_PATH):
    _CUSTOM_FONT_PATHS["demi"] = _DEFAULT_FONT_PATH
    _CUSTOM_FONT_PATHS["medium"] = _DEFAULT_FONT_PATH


def register_fonts(demi_path: Optional[str], medium_path: Optional[str]) -> None:
    """Register actual TTF/OTF font files for accurate width measurement."""
    global _CUSTOM_FONT_PATHS
    if demi_path and os.path.isfile(demi_path):
        _CUSTOM_FONT_PATHS["demi"] = demi_path
    if medium_path and os.path.isfile(medium_path):
        _CUSTOM_FONT_PATHS["medium"] = medium_path


def _get_measure_font(weight: str, size_pt: int) -> ImageFont.FreeTypeFont:
    """Get a PIL font object for measuring text width at a given size."""
    path = _CUSTOM_FONT_PATHS.get(weight)
    if path:
        try:
            return ImageFont.truetype(path, size_pt * 4)  # oversample for accuracy
        except Exception:
            pass
    try:
        return ImageFont.truetype(FALLBACK_FONT_REGULAR, size_pt * 4)
    except Exception:
        return ImageFont.load_default()


def _measure_width_in(text: str, weight: str, size_pt: float) -> float:
    """Measure rendered text width in inches at a given point size."""
    font = _get_measure_font(weight, max(1, int(size_pt)))
    bbox = font.getbbox(text)
    width_px = bbox[2] - bbox[0]
    # font was rendered at size_pt*4 "pixels"; 1pt = 1/72in
    # width_px is in those oversampled pixel units == points * 4
    width_pt = width_px / 4
    return width_pt / 72.0


def fit_font_size(text: str, weight: str, max_pt: float, min_pt: float, max_width_in: float) -> float:
    """Binary-search the largest font size (>= min_pt) at which `text`
    fits within `max_width_in` on a single line. Returns min_pt if even
    min_pt overflows (caller should then truncate or accept overflow)."""
    if not text:
        return max_pt
    lo, hi = min_pt, max_pt
    best = min_pt
    for _ in range(20):
        mid = (lo + hi) / 2
        w = _measure_width_in(text, weight, mid)
        if w <= max_width_in:
            best = mid
            lo = mid
        else:
            hi = mid
        if hi - lo < 0.25:
            break
    return round(best, 1)


def wrap_text_to_width(text: str, weight: str, size_pt: float, max_width_in: float, max_lines: int = 2) -> Optional[list[str]]:
    """Greedy word-wrap `text` into at most `max_lines` lines that each fit
    max_width_in at size_pt. Returns None if it cannot be wrapped within
    max_lines (caller should fall back to merging/shrinking)."""
    words = text.split()
    if not words:
        return [""]
    lines: list[str] = []
    current = ""
    for word in words:
        trial = (current + " " + word).strip()
        if _measure_width_in(trial, weight, size_pt) <= max_width_in:
            current = trial
        else:
            if current:
                lines.append(current)
            current = word
            if len(lines) >= max_lines:
                return None
    if current:
        lines.append(current)
    if len(lines) > max_lines:
        return None
    return lines


_MINOR_WORDS = {
    "a", "an", "the", "of", "and", "or", "for", "to", "in", "on", "at",
    "by", "with", "from", "as", "nor", "but", "is",
}


def smart_title_case(text: str) -> str:
    """Title-case that:
    - does not capitalize the letter after an apostrophe (Hon'ble, not Hon'Ble)
    - keeps minor words (of, and, for, ...) lowercase unless they are the
      first word
    - preserves words that are already all-uppercase (acronyms like IT, M, C)
    """
    if not text:
        return text
    words = text.split(" ")
    out_words = []
    for i, word in enumerate(words):
        # Split off trailing punctuation (commas etc.) to test the core word
        core = word.strip(",.;:")
        if core.isupper() and len(core) <= 4 and core.isalpha():
            # likely an acronym/initial (e.g. "IT", "M", "C", "MP") -> keep as-is
            out_words.append(word)
            continue
        lower_core = core.lower()
        if lower_core in _MINOR_WORDS and i != 0:
            out_words.append(word.lower())
            continue
        # Capitalize only the first alphabetic character; leave the rest
        # (including anything after an apostrophe) as lowercase.
        chars = list(word.lower())
        for j, ch in enumerate(chars):
            if ch.isalpha():
                chars[j] = ch.upper()
                break
        out_words.append("".join(chars))
    return " ".join(out_words)


@dataclass
class Dignitary:
    name: str
    title: str = ""
    company: str = ""


# ---------------------------------------------------------------------------
# Slide building
# ---------------------------------------------------------------------------

def _set_run(run, text, font_name, size_pt, bold=False, color=NAME_COLOR, caps=False):
    run.text = text.upper() if caps else text
    run.font.name = font_name
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.color.rgb = color
    # Ensure east-asian / complex-script font fields also set, for full PPT compatibility
    rPr = run._r.get_or_add_rPr()
    for tag in ("latin", "ea", "cs"):
        el = rPr.find(qn(f"a:{tag}"))
        if el is None:
            el = rPr.makeelement(qn(f"a:{tag}"), {})
            rPr.append(el)
        el.set("typeface", font_name)


def _add_textbox(slide, left_in, top_in, width_in, height_in, rotation=0):
    box = slide.shapes.add_textbox(Inches(left_in), Inches(top_in), Inches(width_in), Inches(height_in))
    box.rotation = rotation
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    return box


def _build_title_company_lines(title: str, company: str, max_width_in: float,
                               max_total_lines: int = 2,
                               title_max_pt: float = TITLE_MAX_PT,
                               title_min_pt: float = TITLE_MIN_PT):
    """Decide layout for title/company per spec:
    - Try title on its own line, company on its own line, each within
      max_total_lines budget combined (default 2 lines total: 1 for title,
      1 for company) at TITLE_MAX_PT, shrinking down to TITLE_MIN_PT.
    - If either overflows even at min size, OR if either would need to
      wrap (i.e. doesn't fit on one line) -> merge into single
      comma-separated line(s) instead.
    Returns (lines: list[str], font_size: float, tight_after_index: set[int])
    tight_after_index marks which line indices should use the TIGHT gap
    to the next line (vs the normal gap before this block).
    """
    title = smart_title_case((title or "").strip())
    company = smart_title_case((company or "").strip())

    if not title and not company:
        return [], title_max_pt
    if title and not company:
        size = fit_font_size(title, "medium", title_max_pt, title_min_pt, max_width_in)
        if _measure_width_in(title, "medium", size) <= max_width_in:
            return [title], size
        wrapped = wrap_text_to_width(title, "medium", title_min_pt, max_width_in, max_lines=2)
        return (wrapped or [title]), title_min_pt
    if company and not title:
        size = fit_font_size(company, "medium", title_max_pt, title_min_pt, max_width_in)
        if _measure_width_in(company, "medium", size) <= max_width_in:
            return [company], size
        wrapped = wrap_text_to_width(company, "medium", title_min_pt, max_width_in, max_lines=2)
        return (wrapped or [company]), title_min_pt

    # Both present: try stacked (own line each) at decreasing size
    for size in [title_max_pt - i * 1.0 for i in range(int((title_max_pt - title_min_pt)) + 1)]:
        title_fits = _measure_width_in(title, "medium", size) <= max_width_in
        company_fits = _measure_width_in(company, "medium", size) <= max_width_in
        if title_fits and company_fits:
            return [title, company], size

    # Stacking failed even at min size -> merge into comma-separated line(s)
    merged = f"{title}, {company}"
    size = fit_font_size(merged, "medium", title_max_pt, title_min_pt, max_width_in)
    if _measure_width_in(merged, "medium", size) <= max_width_in:
        return [merged], size
    wrapped = wrap_text_to_width(merged, "medium", title_min_pt, max_width_in, max_lines=2)
    return (wrapped or [merged]), title_min_pt


def _render_half(slide, dignitary: Dignitary, top_in: float, rotation: int,
                 slide_w: float = SLIDE_W_IN, half_h: float = HALF_H_IN,
                 scale: float = 1.0):
    """Render one half (top or bottom) of the tent card."""
    # Scale all measurements proportionally from A4 defaults
    textbox_w      = TEXTBOX_W_IN        * scale
    margin_x       = (slide_w - textbox_w) / 2
    name_max_pt    = NAME_MAX_PT         * scale
    name_min_pt    = NAME_MIN_PT         * scale
    title_max_pt   = TITLE_MAX_PT        * scale
    title_min_pt   = TITLE_MIN_PT        * scale
    name_title_gap = NAME_TITLE_GAP_IN   * scale
    tc_gap         = TITLE_COMPANY_GAP_IN * scale
    line_spacing   = Pt(50 * scale)

    max_width_in = textbox_w

    name_text = dignitary.name.strip()
    name_size = fit_font_size(name_text, "demi", name_max_pt, name_min_pt, max_width_in)

    title_lines, title_size = _build_title_company_lines(
        dignitary.title, dignitary.company, max_width_in,
        title_max_pt=title_max_pt, title_min_pt=title_min_pt,
    )

    def line_h_in(pt_size):
        return (pt_size * 1.15) / 72.0

    name_h = line_h_in(name_size) * 1.05
    title_block_h = 0.0
    if title_lines:
        if len(title_lines) == 1:
            title_block_h = line_h_in(title_size)
        else:
            title_block_h = line_h_in(title_size) + tc_gap + line_h_in(title_size)

    total_h = name_h + (name_title_gap if title_lines else 0) + title_block_h

    # Use the two-line layout as the reference height for vertical centering
    # so single-line boards look the same as two-line ones.
    two_line_ref_h = name_h + name_title_gap + (2 * line_h_in(title_max_pt) + tc_gap)
    centering_h = max(total_h, two_line_ref_h)
    # d = distance from the fold line to the nearest edge of the name box.
    # Using the same d in both halves makes the layout perfectly symmetric:
    # the name sits at the same distance from the fold whether you're looking
    # at the front or the back of the folded card.
    d = (half_h - centering_h) / 2

    if rotation == 180:
        # Top half: name bottom edge sits at distance d above the fold line.
        name_y  = top_in + half_h - d - name_h
        title_y = name_y - name_title_gap - title_block_h   # title above name
    else:
        # Bottom half: name top edge sits at distance d below the fold line.
        name_y  = top_in + d
        title_y = name_y + name_h + name_title_gap          # title below name

    # --- Name textbox ---
    name_box = _add_textbox(slide, margin_x, name_y, max_width_in, name_h, rotation=rotation)
    p = name_box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    _set_run(run, name_text, FONT_NAME_BOLD, name_size, bold=True, color=NAME_COLOR, caps=True)

    # --- Title/Company textbox(es) ---
    if title_lines:
        title_box = _add_textbox(slide, margin_x, title_y, max_width_in, title_block_h, rotation=rotation)
        tf = title_box.text_frame
        for i, line in enumerate(title_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = PP_ALIGN.CENTER
            p.space_before = Pt(0)
            p.space_after  = Pt(0)
            p.line_spacing = line_spacing
            run = p.add_run()
            _set_run(run, line, FONT_NAME_MEDIUM, title_size, bold=False, color=TITLE_COLOR, caps=False)


# ---------------------------------------------------------------------------
# Paper size presets — all measurements scale proportionally from A4
# ---------------------------------------------------------------------------

PAPER_SIZES = {
    'A4 Landscape': {'w_in': 11.69, 'h_in': 8.27},
    'A5 Landscape': {'w_in':  8.27, 'h_in': 5.83},
}
_A4_W = 11.69  # reference width; all scaling is relative to this


def build_presentation(dignitaries: list[Dignitary], paper_size: str = 'A4 Landscape') -> Presentation:
    size   = PAPER_SIZES.get(paper_size, PAPER_SIZES['A4 Landscape'])
    sw     = size['w_in']
    sh     = size['h_in']
    scale  = sw / _A4_W          # linear scale factor vs A4

    prs = Presentation()
    prs.slide_width  = Inches(sw)
    prs.slide_height = Inches(sh)
    blank_layout = prs.slide_layouts[6]

    half_h = sh / 2

    for dig in dignitaries:
        slide = prs.slides.add_slide(blank_layout)

        # Faint horizontal fold-guide line across the middle (matches sample)
        line = slide.shapes.add_connector(
            1,
            Inches(0.3 * scale), Inches(half_h),
            Inches(sw - 0.3 * scale), Inches(half_h)
        )
        line.line.color.rgb = RGBColor(0xE5, 0xE5, 0xE5)
        line.line.width = Pt(0.25)

        # Top half: rotated 180
        _render_half(slide, dig, top_in=0.0,   rotation=180, slide_w=sw, half_h=half_h, scale=scale)
        # Bottom half: upright
        _render_half(slide, dig, top_in=half_h, rotation=0,   slide_w=sw, half_h=half_h, scale=scale)

    return prs


# ---------------------------------------------------------------------------
# Font embedding
# ---------------------------------------------------------------------------

def _obfuscate_font(font_data: bytes, guid: str) -> bytes:
    """
    Obfuscate font bytes per OOXML spec (ECMA-376 Part 1 §15.2.12.2).

    The first 32 bytes of the font are XOR'd with the 16-byte key derived
    from the GUID. GUID components 1–3 are interpreted little-endian
    (Windows COM convention); components 4–5 are big-endian.
    The 16-byte key is applied twice to cover all 32 bytes.
    """
    g = guid.strip('{}').replace('-', '')
    c1 = bytes(reversed(bytes.fromhex(g[0:8])))   # 4 bytes LE
    c2 = bytes(reversed(bytes.fromhex(g[8:12])))  # 2 bytes LE
    c3 = bytes(reversed(bytes.fromhex(g[12:16]))) # 2 bytes LE
    c4 = bytes.fromhex(g[16:20])                  # 2 bytes BE
    c5 = bytes.fromhex(g[20:32])                  # 6 bytes BE
    key = c1 + c2 + c3 + c4 + c5                  # 16 bytes total
    result = bytearray(font_data)
    for i in range(min(32, len(result))):
        result[i] ^= key[i % 16]
    return bytes(result)


def embed_font_in_pptx(pptx_bytes: bytes, font_path: str, font_name: str) -> bytes:
    """
    Embed a TrueType/OpenType font file directly into PPTX bytes.

    Uses targeted string insertion rather than XML re-parsing, so the
    original XML structure, namespace declarations, and encoding headers
    are preserved byte-for-byte. Only the minimum necessary additions
    are made to register and reference the embedded font.

    The font is obfuscated per the OOXML spec (ECMA-376 §15.2.12.2)
    so it travels with the file and renders correctly on any machine,
    regardless of what fonts are locally installed.
    """
    with open(font_path, 'rb') as fh:
        font_data = fh.read()

    guid           = '{' + str(uuid.uuid4()).upper() + '}'
    rel_id         = guid
    font_part_name = 'ppt/fonts/font1.fntdata'
    obfuscated     = _obfuscate_font(font_data, guid)

    in_buf  = io.BytesIO(pptx_bytes)
    out_buf = io.BytesIO()

    with zipfile.ZipFile(in_buf, 'r') as zin:
        with zipfile.ZipFile(out_buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.infolist():
                data = zin.read(item.filename)
                text = None  # only decode when needed

                if item.filename == '[Content_Types].xml':
                    text = data.decode('utf-8')
                    if 'fntdata' not in text:
                        text = text.replace(
                            '</Types>',
                            '<Default Extension="fntdata" ContentType="application/x-fontdata"/>'
                            '</Types>'
                        )

                elif item.filename == 'ppt/_rels/presentation.xml.rels':
                    text = data.decode('utf-8')
                    text = text.replace(
                        '</Relationships>',
                        f'<Relationship Id="{rel_id}" '
                        f'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/font" '
                        f'Target="fonts/font1.fntdata"/>'
                        '</Relationships>'
                    )

                elif item.filename == 'ppt/presentation.xml':
                    text = data.decode('utf-8')
                    font_block = (
                        f'<p:embeddedFontLst>'
                        f'<p:embeddedFont>'
                        f'<p:font typeface="{font_name}" pitchFamily="0" charset="0"/>'
                        f'<p:regular r:id="{rel_id}"/>'
                        f'</p:embeddedFont>'
                        f'</p:embeddedFontLst>'
                    )
                    # Per the OOXML schema, embeddedFontLst must come after
                    # notesSz and before defaultTextStyle. notesSz is always
                    # self-closing (<p:notesSz .../>) so we anchor on the
                    # element that immediately follows it instead.
                    if '<p:defaultTextStyle>' in text:
                        text = text.replace(
                            '<p:defaultTextStyle>',
                            font_block + '<p:defaultTextStyle>',
                            1
                        )
                    else:
                        # Fallback for unusual presentation structures
                        import re
                        text, n = re.subn(
                            r'(<p:notesSz\b[^>]*/\s*>)',
                            r'\1' + font_block,
                            text, count=1
                        )
                        if not n:
                            text = text.replace(
                                '</p:presentation>',
                                font_block + '</p:presentation>',
                                1
                            )

                zout.writestr(item, text.encode('utf-8') if text is not None else data)

            # Add the obfuscated font binary as a new part
            zout.writestr(font_part_name, obfuscated)

    out_buf.seek(0)
    return out_buf.getvalue()

